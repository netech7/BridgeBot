import io
import fitz
import docx
import pandas as pd
from pptx import Presentation
import pytesseract
from PIL import Image
from langdetect import detect
from langchain.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Set Tesseract path (Windows)
pytesseract.pytesseract.tesseract_cmd = r"C:/Program Files/Tesseract-OCR/tesseract.exe"

# Enable multilingual OCR (supports Hindi, Marathi, Japanese, etc.)
TESSERACT_LANGS = "eng+hin+mar+jpn"

def extract_text_from_file(file):
    name = file.name

    if name.endswith(".pdf"):
        doc = fitz.open(stream=file.read(), filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        return text, name

    elif name.endswith(".docx"):
        doc = docx.Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text, name

    elif name.endswith(".xlsx"):
        df = pd.read_excel(file, sheet_name=None)
        combined_text = ""
        for sheet_name, sheet_data in df.items():
            combined_text += f"Sheet: {sheet_name}\n{sheet_data.to_string(index=False)}\n\n"
        return combined_text, name

    elif name.endswith(".pptx"):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text, name

    elif name.lower().endswith((".png", ".jpg", ".jpeg")):
        image = Image.open(io.BytesIO(file.read()))
        text = pytesseract.image_to_string(image, lang=TESSERACT_LANGS)
        return text, name

    else:
        return "", name

def detect_language(text):
    try:
        return detect(text)
    except:
        return "en"

def process_documents(uploaded_files, embedding_model):
    all_docs = []

    for file in uploaded_files:
        filename = file.name.lower()

        if filename.endswith(".pdf"):
            doc = fitz.open(stream=file.read(), filetype="pdf")
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text()

                image_text = ""
                for img_index, img in enumerate(page.get_images(full=True)):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image = Image.open(io.BytesIO(image_bytes))
                        image_text += pytesseract.image_to_string(image, lang=TESSERACT_LANGS)
                    except Exception as e:
                        print(f"OCR failed on image in PDF page {page_num}: {e}")

                combined_text = text + "\n" + image_text
                lang = detect_language(combined_text)
                if combined_text.strip():
                    all_docs.append(Document(
                        page_content=combined_text,
                        metadata={"source": file.name, "page": page_num, "language": lang}
                    ))

        elif filename.endswith(".docx"):
            doc = docx.Document(file)
            full_text = "\n".join([para.text for para in doc.paragraphs])
            image_text = ""
            for rel in doc.part._rels:
                rel = doc.part._rels[rel]
                if "image" in rel.target_ref:
                    try:
                        img_data = rel.target_part.blob
                        img = Image.open(io.BytesIO(img_data))
                        image_text += pytesseract.image_to_string(img, lang=TESSERACT_LANGS)
                    except Exception as e:
                        print(f"OCR failed in DOCX: {e}")
            combined_text = full_text + "\n" + image_text
            lang = detect_language(combined_text)
            all_docs.append(Document(
                page_content=combined_text,
                metadata={"source": file.name, "page": 1, "language": lang}
            ))

        elif filename.endswith(".xlsx"):
            df = pd.read_excel(file)
            combined_text = df.to_string()
            lang = detect_language(combined_text)
            all_docs.append(Document(
                page_content=combined_text,
                metadata={"source": file.name, "page": 1, "language": lang}
            ))

        elif filename.endswith(".pptx"):
            prs = Presentation(file)
            for slide_num, slide in enumerate(prs.slides, start=1):
                text_runs = []
                image_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
                    if shape.shape_type == 13:
                        try:
                            image = shape.image
                            img_bytes = image.blob
                            img = Image.open(io.BytesIO(img_bytes))
                            image_text += pytesseract.image_to_string(img, lang=TESSERACT_LANGS)
                        except Exception as e:
                            print(f"OCR failed on image in slide {slide_num}: {e}")
                combined_text = "\n".join(text_runs) + "\n" + image_text
                lang = detect_language(combined_text)
                if combined_text.strip():
                    all_docs.append(Document(
                        page_content=combined_text,
                        metadata={"source": file.name, "page": slide_num, "language": lang}
                    ))

        elif filename.endswith((".png", ".jpg", ".jpeg")):
            try:
                image = Image.open(file)
                image_text = pytesseract.image_to_string(image, lang=TESSERACT_LANGS)
                lang = detect_language(image_text)
                all_docs.append(Document(
                    page_content=image_text,
                    metadata={"source": file.name, "page": 1, "language": lang}
                ))
            except Exception as e:
                print(f"Failed OCR on standalone image: {e}")

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=100,
        separators=["\n\n", "\n", ".", " "]
    )
    chunked_docs = text_splitter.split_documents(all_docs)
    return chunked_docs

